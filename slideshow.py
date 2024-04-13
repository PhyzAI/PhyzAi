"""
Implements helper Slide and Slideshow classes
"""
__author__ = "Alexandra"
__version__ = "0.1.0"
__license__ = "NA"

import os
from pptx import Presentation
from PIL import Image

class SlideShow:
    slides = []
    folderPathForPresentation = ""
    slideImagePaths = []
    prs = None


    def __init__(self, p):
        self.folderPathForPresentation = p
        self.slideImagePaths = []
        self.slides = []
        self.prs = None
        print("Trying to find %s" % p)
        # Get Power Point File from folderPathForPresentation
        try:
            # Get the first power point file in specified directory
            for file in os.listdir(self.folderPathForPresentation):
                
                if ".pptx" in file:
                    print(self.folderPathForPresentation + "\\" + file)
                    self.prs = Presentation(".\\" + self.folderPathForPresentation + "\\" + file)
                    break
            # raise Exception('couldnt find file')
        except Exception:
            print("failed to open pptx in %s" % self.folderPathForPresentation)


        # Get Slide Files paths from folderPathForPresentation
        try:
            slideNumber = 0
            while True:
                foundSlide = False
                for file in os.listdir(self.folderPathForPresentation):
                    if ("%i.png" % slideNumber) in file or ("%i.jpg" % slideNumber) in file or ("%i.jpeg" % slideNumber) in file:
                        self.slideImagePaths.append(file)
                        foundSlide = True
                        break
                slideNumber = slideNumber + 1
                if not foundSlide:
                    # print("Found %i slides" % slideNumber)
                    break
                
        except Exception:
            print("failed to find any slides in %s" % self.folderPathForPresentation)

        # Pull Notes from Slides
        slideNumber = 0
        for s in self.prs.slides:
            if slideNumber >= len(self.prs.slides):
                break
            noteText = s.notes_slide.notes_text_frame.text

            # do some processing on the slide's note text
            noteText = noteText.replace("–", ",") # Replace EM dashes with commas
            noteText = noteText.replace("\"", "") # Remove double quotes
            noteText = noteText.replace("\u000b", "") # Remove u+000b chars
            noteText = noteText.replace("`", "") # Remove backticks chars

            # Add this slide to our list of slides in this slideshow
            self.slides.append(Slide(noteText, self.slideImagePaths[slideNumber]))
            slideNumber = slideNumber + 1


# Class for an individual slide
class Slide:
    notes = ""
    imagePath = ""

    def __init__(self, n, i):
        self.notes = n
        self.imagePath = i
    
        
