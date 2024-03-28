"""
A short script to get all slide notes from a pptx slideshow and export it to an XML file.

Call me using format:
python exportNotesFromPPTX.py [path to pptx presentation]

Exports to a xml file with the following structure with the same filename as the pptx file specified:
<notes>
    <note slideNumber="1" text="This is slide 1/>
    <note slideNumber="2" text="This is slide 2/>
    ...
</notes>

"""
__author__ = "Alexandra"
__version__ = "0.1.0"
__license__ = "NA"

import sys
from pptx import Presentation
from xml.dom import minidom 

# Variable Dec
Notes = []
slideshowName = ""
noteFile = minidom.Document()

# Check for proper command format and return format if wrong
if len(sys.argv) != 2:
    print("Expected Command:")
    print("python exportNotesFromPPTX.py [path to pptx presentation]")
    exit()

# Initialize Slideshow
try:
    prs = Presentation(sys.argv[1])
except:
    print("Could not open %s" % sys.argv[1])
    exit()

print("Getting Notes from %s" % sys.argv[1])

# Get info Needed from Slides
slideshowName = sys.argv[1].replace('.pptx', '').replace('.\\', '')
for slide in prs.slides:
    noteText = slide.notes_slide.notes_text_frame.text

    # do some processing on the slide's text
    noteText = noteText.replace("–", ",") # Replace EM dashes with commas
    noteText = noteText.replace("\"", "") # Remove double quotes
    noteText = noteText.replace("\u000b", "") # Remove u+000b chars
    noteText = noteText.replace("`", "") # Remove backticks chars

    Notes.append(noteText)


xml = noteFile.createElement('notes')  
noteFile.appendChild(xml)

# Add Notes to XML File
slideNumber = 1
for note in Notes:
    noteElement = noteFile.createElement("note")
    noteElement.setAttribute("slideNumber", "%s" % slideNumber)
    
    noteElement.setAttribute("text", note)

    xml.appendChild(noteElement)

    slideNumber = slideNumber+1
  

# Export that cute XML!
xmlOutputStr = noteFile.toprettyxml(indent ="\t")

outputPath = "%s.xml" % slideshowName
  
with open(outputPath, "w") as f: 
    f.write(xmlOutputStr)

print("Saved notes as %s" % outputPath)